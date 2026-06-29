# service/document_parser.py
"""文档解析服务"""
import os
import re
import uuid

def parse_document(file_path: str, file_type: str) -> str:
    """解析文档为纯文本"""
    try:
        if file_type == "pdf":
            try:
                import pdfplumber
                text = ""
                with pdfplumber.open(file_path) as pdf:
                    for page in pdf.pages:
                        page_text = page.extract_text()
                        if page_text:
                            text += page_text + "\n"
                return text.strip()
            except ImportError:
                try:
                    from PyPDF2 import PdfReader
                    reader = PdfReader(file_path)
                    text = ""
                    for page in reader.pages:
                        text += page.extract_text() or ""
                    return text.strip()
                except ImportError:
                    return f"[PDF解析失败: 请安装 pdfplumber 或 PyPDF2]"

        elif file_type == "docx":
            try:
                import docx
                doc = docx.Document(file_path)
                paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
                return "\n".join(paragraphs)
            except ImportError:
                return f"[DOCX解析失败: 请安装 python-docx]"

        elif file_type in ["txt", "md"]:
            with open(file_path, "r", encoding="utf-8") as f:
                return f.read()

        elif file_type in ["pptx"]:
            try:
                from pptx import Presentation
                prs = Presentation(file_path)
                text = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            text.append(shape.text.strip())
                return "\n".join(text)
            except ImportError:
                return f"[PPTX解析失败: 请安装 python-pptx]"

        elif file_type in ["xlsx"]:
            try:
                import openpyxl
                wb = openpyxl.load_workbook(file_path)
                text = []
                for sheet in wb.worksheets:
                    for row in sheet.iter_rows(values_only=True):
                        row_text = " ".join([str(cell) for cell in row if cell is not None])
                        if row_text.strip():
                            text.append(row_text)
                return "\n".join(text)
            except ImportError:
                return f"[XLSX解析失败: 请安装 openpyxl]"

        elif file_type in ["png", "jpg", "jpeg", "gif", "bmp", "webp"]:
            return f"[图片文件: {os.path.basename(file_path)}，暂不支持OCR文本提取]"

        else:
            return f"[暂不支持解析 {file_type} 格式，文件路径: {file_path}]"

    except Exception as e:
        raise Exception(f"文档解析失败: {str(e)}")


def split_by_headings(text: str, chunk_size: int = 500, overlap: int = 50) -> list:
    """基于标题层级智能切片"""
    heading_patterns = [
        r'(?:^|\n)#{1,6}\s+[^\n]+',
        r'(?:^|\n)第[一二三四五六七八九十\d]+章\s*[^\n]*',
        r'(?:^|\n)(?:\d+\.)+\s+[^\n]+',
        r'(?:^|\n)\d+\s+[^\n]+',
        r'(?:^|\n)[一二三四五六七八九十]\s*[、.\s]+[^\n]+',
    ]

    combined_pattern = '|'.join(f'({p})' for p in heading_patterns)
    headings = list(re.finditer(combined_pattern, text, re.MULTILINE))

    chunks = []
    if len(headings) == 0:
        return split_fixed(text, chunk_size, overlap)

    for i, match in enumerate(headings):
        start = match.start()
        end = headings[i + 1].start() if i + 1 < len(headings) else len(text)
        content = text[start:end].strip()

        if len(content) > 10:
            heading_text = match.group().strip()
            chunks.append({
                "chunk_index": i,
                "heading_path": heading_text[:100],
                "heading_level": heading_text.count('#') if '#' in heading_text else 1,
                "content": content[:2000],
                "content_length": len(content),
                "page_number": None
            })

    return chunks


def split_fixed(text: str, chunk_size: int = 500, overlap: int = 50) -> list:
    """固定字符数切片"""
    chunks = []
    start = 0
    index = 0
    text_len = len(text)

    while start < text_len:
        end = min(start + chunk_size, text_len)

        if end < text_len:
            for sep in ['。', '；', '\n', ' ']:
                pos = text.rfind(sep, start, end)
                if pos > start + chunk_size // 2:
                    end = pos + len(sep)
                    break

        content = text[start:end].strip()
        if len(content) > 10:
            chunks.append({
                "chunk_index": index,
                "heading_path": f"片段{index+1}",
                "heading_level": 1,
                "content": content,
                "content_length": len(content),
                "page_number": None
            })
            index += 1

        next_start = end - overlap
        if next_start <= start:
            next_start = end
        start = next_start

    return chunks